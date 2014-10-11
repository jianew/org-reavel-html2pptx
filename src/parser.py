#!/usr/bin/env python

from pptx.dml.color import RGBColor

from pptx import Presentation
from pptx.util import Inches, Px
from lxml import etree
from reavelparser import clear_tag as ct

class UnknowenChildException(Exception):
    
    def __init__(self,ptag, ctag):
        Exception.__init__(self)
        self.ptag = ptag
        self.ctag = ctag


class Etree2pptx:

    def __init__(self):
        pass

    def headsection2pptx(self, prs, ein):
        print "head section processed"
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        subts = body_shape.textframe
        
        ein[0].tag = ct(ein[0].tag)

        assert ein[0].tag == "h1"

        title.text = ein[0].text
        
        for head in ein[1:]:
            if len(head.getchildren()) > 0:
                i = head[0]
                assert i.tag.find("a") != -1
                p = subts.add_paragraph()
                p.text = i.text
                
            else:
                p = subts.add_paragraph()
                p.text = head.text
    
    def parse_a2pptx(self, tframe, ein, level):        
        assert ct(ein.tag) == "a"
        p = tframe.add_paragraph()
        run = p.add_run()
        run.text = ein.text
        run.hyperlink.address = ein.get("href")

    def parse_li2pptx(self, tframe,ein,level):
        assert ct(ein.tag) == "li"
        if len(ein.getchildren()) == 0:
            p = tframe.add_paragraph()
            p.text = ein.text
        
        else:
            try:
                for child in ein:
                    if ct(child.tag) == "a":
                        self.parse_a2pptx(tframe, child, level+1)
                    elif ct(child.tag) == "ul":
                        self.parse_ul2pptx(tframe, child, level+1)
                    elif ct(child.tag) == "li":
                        self.parse_li2pptx(tframe, child, level+1)
                    else:
                        raise UnknowenChildException(ein.tag , child.tag)
            except UnknowenChildException,x:
                print "there is no handle of %s\'s chlild %s",(x.ptag,x.ctag)
                    
    def parse_ul2pptx(self, tframe,ein,level):
        assert ct(ein.tag) == "ul"
        for child in ein:
            if ct(child.tag) == "li":
                self.parse_li2pptx(tframe,child,level+1)
            elif ct(child.tag) == "ul":
                self.parse_ul2pptx(tframe,child,level)

    def parse_ol2pptx(self, tframe,ein,level):
        assert ct(ein.tag) == "ol"
        for child in ein:
            if ct(child.tag) == "li":
                self.parse_li2pptx(tframe,child,level)

            elif ct(child.tag) == "ul":
                self.parse_ul2pptx(tframe,child,level+1)

    def is_headtag(self, tagname):
        x = False
        for i in range(1,7):
            if tagname.endswith(str(i)):
                x = True
        return x and tagname.startswith('h')

    def parse_head2pptx(self, title, ein):
        
        assert self.is_headtag(ct(ein.tag))
        
        if ein.text != None:
            title.text = ein.text

    def parse_pre2pptx(self, tframe, ein):
        
        assert ct(ein.tag) == "pre"
        if ein.text == None: return
        p = tframe.add_paragraph()
        run = p.add_run()
        run.text = ein.text
        
        font = run.font
        font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    
    def parse_b2pptx(self, tframe, ein):
        
        assert ct(ein.tag) == "b"
        p = tframe.add_paragraph()
        run = p.add_run()
        run.text = ein.text

        font = run.font
        font.bold = True

    def parse_p2pptx(self , tframe, ein):

        assert ct(ein.tag) == "p"
        if len(ein.getchildren() ) ==0:
            p = tframe.add_paragraph()
            run = p.add_run()
            run.text = ein.text
        else:
            try:
                for child in ein:
                    if ct(child.tag) == "b":
                        self.parse_b2pptx(tframe,child)
                    else:
                        raise  UnknowenChildException(ein.tag,child.tag)
            except   UnknowenChildException,x:
                print "can't handle %s chld %s"%(x.ptag,x.ctag)
    def parse_div2pptx(self, tframe,ein):
        assert ct(ein.tag)== "div"
        assert len(ein.getchildren() ) !=0
        for child in ein:
            if ct(child.tag) == "pre":
                self.parse_pre2pptx(tframe,child)

    def parse_span2pptx(self, tframe, ein):
        
        assert ct(ein.tag) == "span"
        

    def parse_section2pptx(self, prs, ein):

        assert ein.tag.find("section")
        
        section_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(section_slide_layout)
        shapes = slide.shapes
        
        body_shape = shapes.placeholders[1]
        tf = body_shape.textframe

        try:
            for child in ein:

                if self.is_headtag(ct(child.tag)):
                    self.parse_head2pptx(shapes.title, child)

                elif ct(child.tag) == "ul":
                    self.parse_ul2pptx(tf, child ,1)
                
                elif ct(child.tag) == "li":
                    self.parse_li2pptx(tf,child,1)
                
                elif ct(child.tag) == "pre":
                    self.parse_pre2pptx(tf,child)
                    
                elif ct(child.tag) == "ol":
                    self.parse_ol2pptx(tf,child,1)

                elif ct(child.tag) == "section":
                    self.parse_section2pptx(prs,child)
                elif ct(child.tag) == "p":
                    self.parse_p2pptx(tf,child)
                
                elif ct(child.tag) == "div":
                    self.parse_div2pptx(tf, child)
                else:
                    raise UnknowenChildException(ein.tag , child.tag)

        except UnknowenChildException,x:
            print "can't handle %s\'s child %s",(x.ptag,x.ctag)
        
        
